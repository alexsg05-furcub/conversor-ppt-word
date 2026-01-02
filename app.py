import os
import re
from flask import Flask, render_template, request, send_from_directory
from pptx import Presentation
from docx import Document

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "convertidos"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Funciones de tu script
def clean_text(text):
    return re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)

def extract_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text") and shape.text.strip():
        text += clean_text(shape.text) + "\n"
    if shape.shape_type == 19:  # TABLE
        for row in shape.table.rows:
            for cell in row.cells:
                text += clean_text(cell.text) + "\n"
    if shape.shape_type == 6:  # GROUP
        for s in shape.shapes:
            text += extract_text_from_shape(s)
    return text

def pptx_to_docx(pptx_path, docx_path):
    prs = Presentation(pptx_path)
    doc = Document()
    doc.add_heading('Contenido de la presentación', level=1)

    for i, slide in enumerate(prs.slides, start=1):
        doc.add_heading(f'Diapositiva {i}', level=2)
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)
            if text.strip():
                doc.add_paragraph(text)

    doc.save(docx_path)
    print(f"Convertido: {pptx_path} → {docx_path}")

# Rutas web
@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        uploaded_files = request.files.getlist("pptx_files")
        converted_files = []
        for file in uploaded_files:
            if file.filename.lower().endswith(".pptx"):
                pptx_path = os.path.join(UPLOAD_FOLDER, file.filename)
                file.save(pptx_path)
                docx_name = os.path.splitext(file.filename)[0] + ".docx"
                docx_path = os.path.join(OUTPUT_FOLDER, docx_name)
                try:
                    pptx_to_docx(pptx_path, docx_path)
                    converted_files.append(docx_name)
                except Exception as e:
                    print(f"Error al convertir {file.filename}: {e}")
        return render_template("index.html", files=converted_files)
    return render_template("index.html", files=[])

@app.route("/download/<filename>")
def download_file(filename):
    return send_from_directory(OUTPUT_FOLDER, filename, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
